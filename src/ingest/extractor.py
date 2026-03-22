"""
Unified Text Extractor — Edu Nexus
====================================
Single-source text extraction for all supported document formats.
Every extractor returns  List[str]  where each element is one
logical 'page' / 'section' of the document.

Supported formats:
    .pdf   — pdfplumber  (page-level extraction)
    .docx  — python-docx (paragraphs + tables)
    .pptx  — python-pptx (slide shapes + tables + notes)
    .xlsx  — openpyxl    (all sheets, row-level)
    .csv   — stdlib csv  (row-level)
    .txt   — pathlib     (single-page)
    .md    — pathlib     (single-page)
"""

from __future__ import annotations

import csv
import logging
from pathlib import Path
from typing import List

logger = logging.getLogger(__name__)

# ── Supported extensions (single source of truth) ─────────────────────
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".csv", ".txt", ".md"}


# ====================================================================== #
#  PUBLIC API                                                              #
# ====================================================================== #

def extract_text(file_path: Path) -> List[str]:
    """
    Dispatch to the correct extractor based on file extension.

    Returns a list of text strings (one per logical page / section).
    Raises ValueError for unsupported extensions.
    """
    ext = file_path.suffix.lower()

    if ext == ".pdf":
        return _extract_pdf(file_path)
    elif ext == ".docx":
        return _extract_docx(file_path)
    elif ext == ".pptx":
        return _extract_pptx(file_path)
    elif ext == ".xlsx":
        return _extract_xlsx(file_path)
    elif ext == ".csv":
        return _extract_csv(file_path)
    elif ext in (".txt", ".md"):
        return _extract_plaintext(file_path)
    else:
        raise ValueError(
            f"Unsupported file type '{ext}'. "
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )


# ====================================================================== #
#  INDIVIDUAL EXTRACTORS                                                   #
# ====================================================================== #

# ── PDF ────────────────────────────────────────────────────────────────
def _extract_pdf(path: Path) -> List[str]:
    """Extract text page-by-page using pdfplumber."""
    import pdfplumber

    pages: List[str] = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            if text.strip():
                pages.append(text)
    if not pages:
        logger.warning(f"PDF '{path.name}' produced no text (scanned-image PDF?).")
    return pages


# ── DOCX ───────────────────────────────────────────────────────────────
def _extract_docx(path: Path) -> List[str]:
    """Extract paragraphs + table rows from a Word document."""
    from docx import Document

    doc = Document(path)
    parts: List[str] = []

    for p in doc.paragraphs:
        if p.text and p.text.strip():
            parts.append(p.text.strip())

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    text = "\n".join(parts)
    return [text] if text.strip() else []


# ── PPTX ───────────────────────────────────────────────────────────────
def _extract_pptx(path: Path) -> List[str]:
    """
    Extract text from a PowerPoint presentation.

    Each slide becomes one 'page' with content from:
      - Shape text frames (titles, bodies, text boxes)
      - Table cells
      - Speaker notes
    """
    from pptx import Presentation

    prs = Presentation(path)
    pages: List[str] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts: List[str] = []

        # ── Shape text (titles, body, text boxes) ─────────────────────
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_parts.append(text)

            # ── Table cells ───────────────────────────────────────────
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [
                        cell.text.strip()
                        for cell in row.cells
                        if cell.text.strip()
                    ]
                    if cells:
                        slide_parts.append(" | ".join(cells))

        # ── Speaker notes ─────────────────────────────────────────────
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"[Speaker Notes] {notes}")

        if slide_parts:
            page_text = f"--- Slide {slide_num} ---\n" + "\n".join(slide_parts)
            pages.append(page_text)

    if not pages:
        logger.warning(f"PPTX '{path.name}' produced no text.")
    return pages


# ── XLSX ───────────────────────────────────────────────────────────────
def _extract_xlsx(path: Path) -> List[str]:
    """
    Extract text from an Excel workbook.

    Each worksheet becomes one 'page'.  Rows are converted to
    pipe-delimited strings.  The first row (header) is kept intact.
    """
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    pages: List[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_text: List[str] = []

        for row in ws.iter_rows(values_only=True):
            cells = [
                str(cell).strip() for cell in row
                if cell is not None and str(cell).strip()
            ]
            if cells:
                rows_text.append(" | ".join(cells))

        if rows_text:
            page_text = f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows_text)
            pages.append(page_text)

    wb.close()

    if not pages:
        logger.warning(f"XLSX '{path.name}' produced no text.")
    return pages


# ── CSV ────────────────────────────────────────────────────────────────
def _extract_csv(path: Path) -> List[str]:
    """
    Extract text from a CSV file.

    The entire file becomes one 'page'.  Each row is converted to a
    pipe-delimited string for consistency with the XLSX extractor.
    """
    rows_text: List[str] = []

    # Try to auto-detect the dialect
    with open(path, "r", encoding="utf-8", errors="ignore", newline="") as f:
        try:
            dialect = csv.Sniffer().sniff(f.read(4096))
        except csv.Error:
            dialect = csv.excel  # fallback to default comma-separated
        f.seek(0)

        reader = csv.reader(f, dialect)
        for row in reader:
            cells = [cell.strip() for cell in row if cell.strip()]
            if cells:
                rows_text.append(" | ".join(cells))

    if not rows_text:
        logger.warning(f"CSV '{path.name}' produced no text.")
        return []

    return ["\n".join(rows_text)]


# ── Plain text / Markdown ─────────────────────────────────────────────
def _extract_plaintext(path: Path) -> List[str]:
    """Read file as UTF-8 text.  Returns a single-element list."""
    text = path.read_text(encoding="utf-8", errors="ignore")
    return [text] if text.strip() else []
