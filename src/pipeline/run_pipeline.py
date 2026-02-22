import sys
import os
from pathlib import Path
import json
from datetime import datetime
from docx import Document

# Add 'src' to sys.path to allow imports from sibling modules
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

from ingest.ocr import OCR


class UniversalConverter:

    def __init__(self, db_root: Path):
        self.db_root = db_root
        self.raw = db_root / "raw"
        self.intermediate = db_root / "intermediate"
        self.normalized = db_root / "normalized"
        self.metadata_file = db_root / "metadata" / "processing_log.json"

        self.ocr = OCR()

        (self.intermediate / "ocr_text").mkdir(parents=True, exist_ok=True)
        (self.normalized / "docx").mkdir(parents=True, exist_ok=True)
        (self.db_root / "metadata").mkdir(parents=True, exist_ok=True)

    # ---------------- IMAGE PROCESSING ----------------
    def process_images(self):
        folder = self.raw / "images"

        for img in folder.glob("*"):
            print(f"[IMAGE OCR] {img.name}")

            text = self.ocr.image_to_text(str(img))
            self.save_outputs(img.stem, text)

    # ---------------- PDF PROCESSING ----------------
    def process_pdfs(self):
        import fitz  # PyMuPDF

        folder = self.raw / "pdf"

        for pdf in folder.glob("*.pdf"):
            print(f"[PDF] {pdf.name}")
            text = ""

            doc = fitz.open(pdf)
            for page in doc:
                page_text = page.get_text()
                text += page_text

            if len(text.strip()) < 50:
                print(" → OCR fallback")
                # Convert pages to images
                for i, page in enumerate(doc):
                    pix = page.get_pixmap()
                    img_path = self.intermediate / f"{pdf.stem}_{i}.png"
                    pix.save(img_path)

                    text += self.ocr.image_to_text(str(img_path))

            self.save_outputs(pdf.stem, text)

    # ---------------- PPT PROCESSING ----------------
    def process_ppts(self):
        from pptx import Presentation

        folder = self.raw / "ppt"

        for ppt in folder.glob("*.pptx"):
            print(f"[PPT] {ppt.name}")
            text = ""

            prs = Presentation(ppt)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

            self.save_outputs(ppt.stem, text)

    # ---------------- SAVE OUTPUT ----------------
    def save_outputs(self, name, text):
        # save intermediate
        with open(self.intermediate / "ocr_text" / f"{name}.txt", "w", encoding="utf-8") as f:
            f.write(text)

        # save docx
        doc = Document()
        doc.add_paragraph(text)
        doc.save(self.normalized / "docx" / f"{name}.docx")

        self.update_metadata(name)

    # ---------------- METADATA ----------------
    def update_metadata(self, filename):
        if self.metadata_file.exists():
            try:
                with open(self.metadata_file, "r", encoding="utf-8") as f:
                    metadata = json.load(f)
            except json.JSONDecodeError:
                metadata = []
        else:
            metadata = []

        metadata.append({
            "file_name": filename,
            "status": "processed",
            "timestamp": datetime.now().isoformat()
        })

        with open(self.metadata_file, "w", encoding="utf-8") as f:
            json.dump(metadata, f, indent=4)

    # ---------------- RUN ALL ----------------
    def run_all(self):
        self.process_images()
        self.process_pdfs()
        self.process_ppts()


if __name__ == "__main__":
    repo_root = Path(__file__).resolve().parents[2]
    db_root = repo_root.parent / "edu_nexus_db"

    pipeline = UniversalConverter(db_root)
    pipeline.run_all()

    print("✅ Universal pipeline completed")
